"""
Knowledge Warehouse - AI Presentation Generator
Converts Word documents (.docx / .odt) to PowerPoint with optional AI beautification.
Developed purely using AI LLM models. Motivation by Knowledge Warehouse.
"""

import streamlit as st
import sqlite3
import secrets
import string
import os
import smtplib
import ssl
import json
import base64
import re
import io
from datetime import datetime, timedelta
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart

from docx import Document
from docx.oxml.ns import qn as wqn

try:
    from odf.opendocument import load as load_odt
    from odf.text import P as OdfP
    ODF_AVAILABLE = True
except ImportError:
    ODF_AVAILABLE = False

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn as pqn

import requests

# ═══════════════════════════════════════════════
# CONFIGURATION — secrets come from Replit Secrets
# ═══════════════════════════════════════════════
DB_FILE = "kw_pptgen.db"
ENCRYPTION_SECRET = os.environ.get("ENCRYPTION_SECRET", "kw-change-this-secret-key")
SUPERUSER_USERNAME = os.environ.get("SUPERUSER_USERNAME", "admin")
SUPERUSER_PASSWORD = os.environ.get("SUPERUSER_PASSWORD", "superadmin123")
SMTP_SERVER = os.environ.get("SMTP_SERVER", "")
SMTP_PORT = int(os.environ.get("SMTP_PORT", "587"))
SMTP_USERNAME = os.environ.get("SMTP_USERNAME", "")
SMTP_PASSWORD = os.environ.get("SMTP_PASSWORD", "")
SENDER_EMAIL = os.environ.get("SENDER_EMAIL", "")
APP_NAME = "Knowledge Warehouse – PPT Generator"
FREE_USAGE_LIMIT = 2

LLM_MODELS = {
    "ChatGPT (GPT-4o)": {"url": "https://api.openai.com/v1/chat/completions",
                          "model": "gpt-4o", "type": "openai"},
    "ChatGPT (GPT-4o Mini)": {"url": "https://api.openai.com/v1/chat/completions",
                               "model": "gpt-4o-mini", "type": "openai"},
    "ChatGPT (GPT-3.5 Turbo)": {"url": "https://api.openai.com/v1/chat/completions",
                                  "model": "gpt-3.5-turbo", "type": "openai"},
    "DeepSeek (V3)": {"url": "https://api.deepseek.com/chat/completions",
                      "model": "deepseek-chat", "type": "openai"},
    "DeepSeek (R1)": {"url": "https://api.deepseek.com/chat/completions",
                      "model": "deepseek-reasoner", "type": "openai"},
    "Zephyr 7B (HuggingFace)": {"url": "https://api-inference.huggingface.co/models/HuggingFaceH4/zephyr-7b-beta",
                                 "model": "HuggingFaceH4/zephyr-7b-beta", "type": "huggingface"},
    "Llama 3.1 70B (Groq)": {"url": "https://api.groq.com/openai/v1/chat/completions",
                              "model": "llama-3.1-70b-versatile", "type": "openai"},
    "Llama 3.3 70B (Groq)": {"url": "https://api.groq.com/openai/v1/chat/completions",
                              "model": "llama-3.3-70b-versatile", "type": "openai"},
    "Nemotron 70B (NVIDIA)": {"url": "https://integrate.api.nvidia.com/v1/chat/completions",
                              "model": "nvidia/llama-3.1-nemotron-70b-instruct", "type": "openai"},
    "Custom OpenAI-Compatible": {"url": "", "model": "", "type": "openai"},
}

# ═══════════════════════════════════════════════
# DATABASE
# ═══════════════════════════════════════════════
def get_db():
    conn = sqlite3.connect(DB_FILE)
    conn.row_factory = sqlite3.Row
    return conn

def init_db():
    conn = get_db()
    conn.execute("""CREATE TABLE IF NOT EXISTS users (
        code TEXT PRIMARY KEY,
        usage_count INTEGER DEFAULT 0,
        is_blocked INTEGER DEFAULT 0,
        access_expiry TEXT,
        encrypted_api_key TEXT,
        api_model TEXT,
        created_at TEXT,
        last_login TEXT
    )""")
    conn.commit()
    conn.close()

def get_user(code):
    conn = get_db()
    row = conn.execute("SELECT * FROM users WHERE code = ?", (code,)).fetchone()
    conn.close()
    return dict(row) if row else None

def create_user(code):
    conn = get_db()
    now = datetime.now().isoformat()
    conn.execute("INSERT INTO users (code, usage_count, is_blocked, access_expiry, created_at) VALUES (?,?,?,?,?)",
                 (code, 0, 0, None, now))
    conn.commit()
    conn.close()

def increment_usage(code):
    conn = get_db()
    now = datetime.now().isoformat()
    conn.execute("UPDATE users SET usage_count = usage_count + 1, last_login = ? WHERE code = ?", (now, code))
    conn.commit()
    conn.close()

def get_all_users():
    conn = get_db()
    rows = conn.execute("SELECT * FROM users ORDER BY created_at DESC").fetchall()
    conn.close()
    return [dict(r) for r in rows]

def unblock_user(code, months):
    conn = get_db()
    expiry = (datetime.now() + timedelta(days=30 * months)).isoformat()
    conn.execute("UPDATE users SET is_blocked = 0, access_expiry = ? WHERE code = ?", (expiry, code))
    conn.commit()
    conn.close()

def block_user(code):
    conn = get_db()
    conn.execute("UPDATE users SET is_blocked = 1, access_expiry = NULL WHERE code = ?", (code,))
    conn.commit()
    conn.close()

def save_user_api(code, encrypted_key, model_name):
    conn = get_db()
    conn.execute("UPDATE users SET encrypted_api_key = ?, api_model = ? WHERE code = ?",
                 (encrypted_key, model_name, code))
    conn.commit()
    conn.close()

# ═══════════════════════════════════════════════
# ENCRYPTION UTILITIES
# ═══════════════════════════════════════════════
def _xor_crypt(data_bytes, key_bytes):
    return bytes(a ^ b for a, b in zip(data_bytes, key_bytes * (len(data_bytes) // len(key_bytes) + 1)))

def encrypt_api_key(api_key):
    return base64.b64encode(_xor_crypt(api_key.encode(), ENCRYPTION_SECRET.encode())).decode()

def decrypt_api_key(enc):
    try:
        return _xor_crypt(base64.b64decode(enc), ENCRYPTION_SECRET.encode()).decode()
    except Exception:
        return ""

# ═══════════════════════════════════════════════
# EMAIL
# ═══════════════════════════════════════════════
def send_code_email(recipient, code):
    if not all([SMTP_SERVER, SMTP_USERNAME, SMTP_PASSWORD, SENDER_EMAIL]):
        return False, "SMTP_NOT_CONFIGURED"
    try:
        msg = MIMEMultipart("alternative")
        msg["Subject"] = f"Your Access Code – {APP_NAME}"
        msg["From"] = SENDER_EMAIL
        msg["To"] = recipient
        txt = (f"Your access code for {APP_NAME} is:\n\n{code}\n\n"
               "Keep this code safe. You need it to log in.\n"
               "— Knowledge Warehouse (www.knwaho.com)")
        htm = (f"<html><body><h2>Your Access Code</h2>"
               f"<p>Code for <b>{APP_NAME}</b>:</p>"
               f"<div style='font-family:monospace;font-size:1.6em;background:#f0f4f8;"
               f"padding:18px;text-align:center;color:#1a3c5e;border-radius:8px;'>{code}</div>"
               f"<p>Keep it safe. — <a href='https://www.knwaho.com'>Knowledge Warehouse</a></p>"
               f"</body></html>")
        msg.attach(MIMEText(txt, "plain"))
        msg.attach(MIMEText(htm, "html"))
        ctx = ssl.create_default_context()
        with smtplib.SMTP(SMTP_SERVER, SMTP_PORT) as s:
            s.starttls(context=ctx)
            s.login(SMTP_USERNAME, SMTP_PASSWORD)
            s.sendmail(SENDER_EMAIL, recipient, msg.as_string())
        return True, "Code sent to your email successfully!"
    except Exception as e:
        return False, f"Email failed ({e}). Your code is shown below — please save it."

# ═══════════════════════════════════════════════
# ACCESS CHECK
# ═══════════════════════════════════════════════
def check_access(code):
    user = get_user(code)
    if not user:
        return False, "Invalid access code."
    now = datetime.now()
    if user["access_expiry"]:
        exp = datetime.fromisoformat(user["access_expiry"])
        if exp > now:
            return True, "Access granted (subscription active)."
        else:
            return False, ("Your subscription has expired. Please contact "
                           "Knowledge Warehouse at www.knwaho.com for renewal.")
    if user["is_blocked"]:
        return False, ("Your access is blocked. Please contact "
                       "Knowledge Warehouse at www.knwaho.com.")
    if user["usage_count"] >= FREE_USAGE_LIMIT:
        return False, (f"You have exhausted your {FREE_USAGE_LIMIT} free uses. "
                       "To continue, please contact Knowledge Warehouse at www.knwaho.com.")
    remaining = FREE_USAGE_LIMIT - user["usage_count"]
    return True, f"Access granted ({remaining} free use{'s' if remaining != 1 else ''} remaining)."

# ═══════════════════════════════════════════════
# WORD DOCUMENT PARSING
# ═══════════════════════════════════════════════
def _is_heading(style_name):
    return style_name and style_name.lower().startswith("heading")

def _heading_level(style_name):
    if not style_name:
        return 999
    m = re.search(r"(\d+)", style_name)
    return int(m.group(1)) if m else 999

def _extract_docx_content(file_bytes):
    doc = Document(io.BytesIO(file_bytes))
    items = []
    para_map = {id(p._element): p for p in doc.paragraphs}
    table_map = {id(t._element): t for t in doc.tables}

    for child in doc.element.body:
        tag = child.tag
        if tag == wqn("w:p"):
            p = para_map.get(id(child))
            if p is None:
                continue
            text = p.text.strip()
            style = p.style.name if p.style else "Normal"
            has_pb = False
            pPr = child.find(wqn("w:pPr"))
            if pPr is not None:
                if pPr.find(wqn("w:sectPr")) is not None:
                    has_pb = True
            for run_el in child.findall(wqn("w:r")):
                for br in run_el.findall(wqn("w:br")):
                    if br.get(wqn("w:type")) == "page":
                        has_pb = True
            if text:
                items.append({"type": "para", "text": text, "style": style, "pb": has_pb})
        elif tag == wqn("w:tbl"):
            t = table_map.get(id(child))
            if t is None:
                continue
            rows_text = []
            for row in t.rows:
                cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                rows_text.append(" | ".join(cells))
            full = "\n".join(rows_text).strip()
            if full:
                items.append({"type": "table", "text": full, "style": "Table", "pb": False})
    return items

def _extract_odt_content(file_bytes):
    if not ODF_AVAILABLE:
        st.error("odfpy is not installed — .odt files are not supported.")
        return []
    tmp = io.BytesIO(file_bytes)
    doc = load_odt(tmp)
    items = []
    for elem in doc.text.getElementsByType(OdfP):
        txt = ""
        for ch in elem.childNodes:
            txt += (ch.data if hasattr(ch, "data") else "")
        txt = txt.strip()
        if txt:
            items.append({"type": "para", "text": txt, "style": "Normal", "pb": False})
    return items

def split_into_slides(items):
    slides, cur = [], []
    for it in items:
        if it["pb"] and cur:
            slides.append(cur)
            cur = []
        cur.append(it)
    if cur:
        slides.append(cur)
    if len(slides) > 1:
        return slides

    slides, cur = [], []
    for it in items:
        if _is_heading(it["style"]) and _heading_level(it["style"]) == 1 and cur:
            slides.append(cur)
            cur = []
        cur.append(it)
    if cur:
        slides.append(cur)
    if len(slides) > 1:
        return slides

    slides, cur = [], []
    for it in items:
        if _is_heading(it["style"]) and _heading_level(it["style"]) <= 2 and cur:
            slides.append(cur)
            cur = []
        cur.append(it)
    if cur:
        slides.append(cur)
    if len(slides) > 1:
        return slides

    slides, cur = [], []
    CHUNK = 6
    count = 0
    for it in items:
        cur.append(it)
        if it["type"] == "para":
            count += 1
        if count >= CHUNK:
            slides.append(cur)
            cur = []
            count = 0
    if cur:
        slides.append(cur)
    return slides if slides else [items]

# ═══════════════════════════════════════════════
# PPTX GENERATION
# ═══════════════════════════════════════════════
def _add_text_frame(slide, left, top, width, height, text_lines, font_size=18,
                    bold_first=False, color=RGBColor(0x33, 0x33, 0x33), alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
        if bold_first and i == 0:
            run.font.bold = True
            run.font.size = Pt(font_size + 6)
    return txBox

def build_pptx(slide_contents):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for idx, items in enumerate(slide_contents):
        slide = prs.slides.add_slide(blank_layout)
        title_lines = []
        body_lines = []
        for it in items:
            if _is_heading(it["style"]) and _heading_level(it["style"]) <= 2:
                title_lines.append(it["text"])
            else:
                for line in it["text"].split("\n"):
                    body_lines.append(line)

        if not title_lines and not body_lines:
            continue

        if title_lines:
            _add_text_frame(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.5),
                            title_lines, font_size=28, bold_first=True,
                            color=RGBColor(0x1a, 0x3c, 0x5e), alignment=PP_ALIGN.LEFT)

        body_top = Inches(2.2) if title_lines else Inches(0.8)
        body_height = Inches(5.0) if title_lines else Inches(6.2)
        if body_lines:
            _add_text_frame(slide, Inches(0.8), body_top, Inches(11.7), body_height,
                            body_lines, font_size=18,
                            color=RGBColor(0x33, 0x33, 0x33), alignment=PP_ALIGN.LEFT)

        txBox = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = str(idx + 1)
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

# ═══════════════════════════════════════════════
# AI BEAUTIFICATION
# ═══════════════════════════════════════════════
def hex_to_rgb(h):
    h = h.lstrip("#")
    if len(h) != 6:
        return None
    try:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        return None

def _extract_json(text):
    m = re.search(r"```(?:json)?\s*([\s\S]*?)```", text)
    if m:
        try:
            return json.loads(m.group(1))
        except Exception:
            pass
    m = re.search(r"\{[\s\S]*\}", text)
    if m:
        try:
            return json.loads(m.group(0))
        except Exception:
            pass
    return None

def call_llm(model_info, api_key, prompt):
    if model_info["type"] == "huggingface":
        headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
        payload = {"inputs": prompt, "parameters": {"temperature": 0.4, "max_new_tokens": 4096,
                                                     "return_full_text": False}}
        r = requests.post(model_info["url"], headers=headers, json=payload, timeout=120)
        r.raise_for_status()
        data = r.json()
        if isinstance(data, list) and len(data) > 0:
            return data[0].get("generated_text", "")
        return str(data)
    else:
        headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
        payload = {"model": model_info["model"], "messages": [
            {"role": "system", "content": "You are a presentation design expert. Return ONLY valid JSON."},
            {"role": "user", "content": prompt}
        ], "temperature": 0.4, "max_tokens": 4096}
        r = requests.post(model_info["url"], headers=headers, json=payload, timeout=120)
        r.raise_for_status()
        return r.json()["choices"][0]["message"]["content"]

def build_beautify_prompt(slide_contents):
    slides_desc = []
    for i, items in enumerate(slide_contents):
        texts = [it["text"][:200] for it in items if it["text"]]
        slides_desc.append(f"Slide {i+1}: {'; '.join(texts)}")
    return f"""You are an expert presentation designer. I have a PowerPoint with {len(slide_contents)} slides.
Suggest a beautiful, professional, cohesive color scheme and per-slide formatting.

Return ONLY valid JSON (no markdown, no extra text) in this exact structure:
{{
  "theme": {{
    "title_font": "Calibri",
    "title_size": 30,
    "title_color": "#1a3c5e",
    "body_font": "Calibri",
    "body_size": 18,
    "body_color": "#333333"
  }},
  "slides": [
    {{
      "slide_index": 0,
      "background_color": "#1a3c5e",
      "title_color": "#FFFFFF",
      "body_color": "#E0E0E0",
      "title_alignment": "center",
      "body_alignment": "center"
    }},
    {{
      "slide_index": 1,
      "background_color": "#FFFFFF",
      "title_color": "#1a3c5e",
      "body_color": "#333333",
      "title_alignment": "left",
      "body_alignment": "left"
    }}
  ]
}}

Rules:
- Use hex colors only (e.g. "#1a3c5e").
- First slide (index 0) should be a title slide with a dark background and light text.
- Content slides should have light backgrounds with dark text.
- Be creative with accent colors but keep it professional.
- You must output entries for ALL {len(slide_contents)} slides.

Slide contents:
{chr(10).join(slides_desc)}"""

def apply_beautification(pptx_bytes, design_json, slide_contents):
    prs = Presentation(io.BytesIO(pptx_bytes))
    theme = design_json.get("theme", {})
    slides_map = {s.get("slide_index", i): s for i, s in enumerate(design_json.get("slides", []))}

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
                 "justify": PP_ALIGN.JUSTIFY}

    for idx, slide in enumerate(prs.slides):
        sd = slides_map.get(idx, {})
        bg_hex = sd.get("background_color")
        if bg_hex:
            rgb = hex_to_rgb(bg_hex)
            if rgb:
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = rgb

        title_color = hex_to_rgb(sd.get("title_color", theme.get("title_color", "#1a3c5e"))) or RGBColor(0x1a, 0x3c, 0x5e)
        body_color = hex_to_rgb(sd.get("body_color", theme.get("body_color", "#333333"))) or RGBColor(0x33, 0x33, 0x33)
        title_font = theme.get("title_font", "Calibri")
        title_size = theme.get("title_size", 30)
        body_font = theme.get("body_font", "Calibri")
        body_size = theme.get("body_size", 18)
        title_align = align_map.get(sd.get("title_alignment", "left"), PP_ALIGN.LEFT)
        body_align = align_map.get(sd.get("body_alignment", "left"), PP_ALIGN.LEFT)

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            is_title = False
            if tf.paragraphs:
                first_text = tf.paragraphs[0].text.strip()
                if idx < len(slide_contents) and slide_contents[idx]:
                    first_item_text = slide_contents[idx][0].get("text", "")
                    if first_text and first_item_text and first_text in first_item_text:
                        is_title = True

            for pi, para in enumerate(tf.paragraphs):
                para.alignment = title_align if (is_title and pi == 0) else body_align
                for run in para.runs:
                    if is_title and pi == 0:
                        run.font.name = title_font
                        run.font.size = Pt(title_size)
                        run.font.color.rgb = title_color
                        run.font.bold = True
                    else:
                        run.font.name = body_font
                        run.font.size = Pt(body_size)
                        run.font.color.rgb = body_color

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

# ═══════════════════════════════════════════════
# CUSTOM CSS
# ═══════════════════════════════════════════════
CUSTOM_CSS = """
<style>
    .kw-header {
        background: linear-gradient(135deg, #0d1b2a 0%, #1b3a5c 50%, #274c77 100%);
        padding: 2.5rem 2rem 2rem 2rem;
        border-radius: 12px;
        margin-bottom: 2rem;
        text-align: center;
    }
    .kw-header h1 {
        color: #e0e7ee;
        font-size: 2.2rem;
        margin-bottom: 0.3rem;
        font-weight: 700;
        letter-spacing: 0.5px;
    }
    .kw-header p.sub {
        color: #8badc4;
        font-size: 1.05rem;
    }
    .kw-disclaimer {
        background: #fff8e1;
        border-left: 4px solid #ffa000;
        padding: 1.2rem 1.5rem;
        border-radius: 6px;
        margin: 1.5rem 0;
        font-size: 0.92rem;
        color: #4a4a4a;
        line-height: 1.65;
    }
    .kw-disclaimer strong {
        color: #1a3c5e;
    }
    .kw-card {
        background: #f8fafc;
        border: 1px solid #dde5ed;
        border-radius: 10px;
        padding: 1.8rem;
        margin: 1rem 0;
    }
    .kw-code-display {
        background: #0d1b2a;
        color: #4cc9f0;
        font-family: 'Courier New', monospace;
        font-size: 1.8rem;
        padding: 1rem 2rem;
        border-radius: 8px;
        text-align: center;
        letter-spacing: 3px;
        margin: 1rem 0;
        user-select: all;
    }
    div[data-testid="stTabs"] button[data-baseweb="tab"] {
        font-size: 1rem;
        padding: 0.6rem 1.5rem;
    }
</style>
"""

# ═══════════════════════════════════════════════
# UI: LANDING PAGE
# ═══════════════════════════════════════════════
def show_landing():
    st.markdown(CUSTOM_CSS, unsafe_allow_html=True)
    st.markdown("""<div class="kw-header">
        <h1>📄 AI Presentation Generator</h1>
        <p class="sub">Convert Word documents to beautiful PowerPoint presentations</p>
    </div>""", unsafe_allow_html=True)

    st.markdown("""<div class="kw-disclaimer">
        <strong>Application by Knowledge Warehouse.</strong> The motivation and enthusiasm behind creating this application goes to <strong>Knowledge Warehouse</strong>. The company is not liable for any data loss, malfunctioning and data leaks though the application has been tested thoroughly. Users at their discretion and utmost care may use the application. Though the output in most cases are accurate but sometimes may not be up to expectation. For demonstration and personal use only. Please don't input personal and sensitive information. Application to be used with utmost prudence. Our application does not store any personal information including their email ids which they are logging in with.<br><br>
        To develop more such personalized applications and to discuss on commercial aspect, kindly get in touch with Knowledge Warehouse on their website @ <strong><a href="https://www.knwaho.com" target="_blank">www.knwaho.com</a></strong>
    </div>""", unsafe_allow_html=True)

    mode = st.radio("Choose an option", ["🔑 I have an access code – Login",
                                          "✉️  I'm new – Register with email"],
                     horizontal=True, label_visibility="collapsed")

    if "🔑" in mode:
        st.markdown('<div class="kw-card">', unsafe_allow_html=True)
        code = st.text_input("Enter your access code", placeholder="e.g. KW-XXXX-XXXX-XXXX", key="login_code")
        if st.button("Login", type="primary", use_container_width=True):
            if not code.strip():
                st.error("Please enter your access code.")
            else:
                ok, msg = check_access(code.strip())
                if ok:
                    increment_usage(code.strip())
                    st.session_state.logged_in = True
                    st.session_state.user_code = code.strip()
                    st.session_state.is_superuser = False
                    st.rerun()
                else:
                    st.error(msg)
        st.markdown('</div>', unsafe_allow_html=True)

    else:
        st.markdown('<div class="kw-card">', unsafe_allow_html=True)
        email = st.text_input("Enter your email address", placeholder="you@example.com", key="reg_email")
        if st.button("Register & Get Code", type="primary", use_container_width=True):
            if not email.strip():
                st.error("Please enter your email address.")
            elif not re.match(r'^[a-zA-Z0-9._%+\-]+@[a-zA-Z0-9.\-]+\.[a-zA-Z]{2,}$', email.strip()):
                st.error("Please enter a valid email address.")
            else:
                code = "KW-" + "-".join(secrets.token_hex(2).upper() for _ in range(3))
                create_user(code)
                ok, msg = send_code_email(email.strip(), code)
                st.session_state.reg_code = code
                st.session_state.reg_email_msg = msg
                st.session_state.reg_email_ok = ok
                st.rerun()

        if st.session_state.get("reg_code"):
            if st.session_state.get("reg_email_ok"):
                st.success("✅ Code sent to your email successfully!")
            else:
                if st.session_state.reg_email_msg == "SMTP_NOT_CONFIGURED":
                    st.error("⚠️ EMAIL SYSTEM IS NOT CONFIGURED. Since the code cannot be emailed, please **COPY** the code below and save it securely. You will need it to log in!")
                else:
                    st.warning(st.session_state.reg_email_msg)
            
            st.markdown(f'<div class="kw-code-display">{st.session_state.reg_code}</div>',
                        unsafe_allow_html=True)
            st.info("⚠️ Please save this code now. It will not be shown again. Your email has been permanently deleted from our system.")
        st.markdown('</div>', unsafe_allow_html=True)

    with st.expander("Superuser Login"):
        su_id = st.text_input("Superuser ID", key="su_id")
        sp = st.text_input("Superuser Password", type="password", key="su_pass")
        if st.button("Login as Superuser"):
            if not su_id or not sp:
                st.error("Please enter both Superuser ID and Password.")
            elif su_id == SUPERUSER_USERNAME and sp == SUPERUSER_PASSWORD:
                st.session_state.logged_in = True
                st.session_state.user_code = "__SUPERUSER__"
                st.session_state.is_superuser = True
                st.rerun()
            else:
                st.error("Invalid Superuser ID or Password.")

# ═══════════════════════════════════════════════
# UI: TAB 1 – GENERATE PRESENTATION
# ═══════════════════════════════════════════════
def show_tab1():
    st.header("📄 Generate Presentation from Word File")
    st.markdown("Upload a **.docx** (Microsoft Word) or **.odt** (OpenOffice) file. "
                "Content per page becomes content per slide. "
                "For best results, use **page breaks** or **Heading 1** styles in your document.")

    uploaded = st.file_uploader("Choose a Word file", type=["docx", "odt"],
                                key="word_upload_tab1")

    if uploaded:
        ext = uploaded.name.rsplit(".", 1)[-1].lower()
        with st.spinner("Parsing document..."):
            try:
                if ext == "docx":
                    items = _extract_docx_content(uploaded.read())
                elif ext == "odt":
                    if not ODF_AVAILABLE:
                        st.error("odfpy library is not installed. .odt files are not supported on this instance.")
                        return
                    items = _extract_odt_content(uploaded.read())
                else:
                    st.error("Unsupported file type. Please upload .docx or .odt files only.")
                    return
            except Exception as e:
                st.error(f"Error reading file: {e}")
                return

        if not items:
            st.warning("No text content found in the document.")
            return

        slides = split_into_slides(items)

        st.success(f"Document parsed: {len(items)} content blocks → {len(slides)} slide(s)")

        with st.expander("Preview slide contents", expanded=False):
            for i, sl in enumerate(slides):
                st.markdown(f"**Slide {i+1}:**")
                for it in sl:
                    prefix = "📋 " if it["type"] == "table" else ""
                    st.text(f"  {prefix}{it['text'][:150]}{'...' if len(it['text'])>150 else ''}")
                st.divider()

        if st.button("⚙️ Generate Presentation", type="primary", use_container_width=True):
            with st.spinner("Building PowerPoint..."):
                try:
                    pptx_bytes = build_pptx(slides)
                    st.session_state.base_pptx = pptx_bytes
                    st.session_state.slide_contents = slides
                    st.session_state.base_filename = uploaded.name.rsplit(".", 1)[0] + ".pptx"
                    st.success("Presentation generated successfully!")
                except Exception as e:
                    st.error(f"Error generating presentation: {e}")

    if st.session_state.get("base_pptx"):
        st.download_button(
            label="⬇️ Download Presentation",
            data=st.session_state.base_pptx,
            file_name=st.session_state.get("base_filename", "presentation.pptx"),
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )

# ═══════════════════════════════════════════════
# UI: TAB 2 – BEAUTIFY WITH AI
# ═══════════════════════════════════════════════
def show_tab2():
    st.header("🎨 Beautify Presentation with AI")
    user = get_user(st.session_state.user_code)
    if not user:
        st.error("User not found.")
        return

    if not st.session_state.get("base_pptx"):
        st.warning("Please generate a presentation in **Tab 1** first, then come here to beautify it.")
        return

    st.markdown('<div class="kw-card">', unsafe_allow_html=True)
    st.subheader("AI Model Configuration")
    model_name = st.selectbox("Select LLM Model", list(LLM_MODELS.keys()),
                               index=0, key="beaut_model")

    selected = LLM_MODELS[model_name]
    custom_url = ""
    custom_model = ""
    if model_name == "Custom OpenAI-Compatible":
        custom_url = st.text_input("API Endpoint URL", placeholder="https://api.example.com/v1/chat/completions")
        custom_model = st.text_input("Model Name", placeholder="my-model")

    stored_key = decrypt_api_key(user["encrypted_api_key"]) if user["encrypted_api_key"] else ""
    stored_model = user["api_model"] or ""

    if stored_key and stored_model == model_name:
        st.info("🔑 A saved API key exists for this model. You can update it below or leave blank to use the saved one.")
        api_key = st.text_input("API Key (leave blank to use saved)", type="password", key="beaut_key")
        if not api_key:
            api_key = stored_key
    else:
        api_key = st.text_input("API Key", type="password", key="beaut_key", placeholder="sk-...")

    if not api_key:
        st.error("Please enter your API key.")
        st.markdown('</div>', unsafe_allow_html=True)
        return

    enc_key = encrypt_api_key(api_key)
    save_user_api(st.session_state.user_code, enc_key, model_name)

    effective_url = custom_url if model_name == "Custom OpenAI-Compatible" else selected["url"]
    effective_model = custom_model if model_name == "Custom OpenAI-Compatible" else selected["model"]
    effective_type = selected["type"]

    st.markdown('</div>', unsafe_allow_html=True)

    if st.button("✨ Beautify Presentation", type="primary", use_container_width=True):
        if not effective_url:
            st.error("Please provide the API endpoint URL for the custom model.")
            return
        if not effective_model:
            st.error("Please provide the model name for the custom model.")
            return

        slides = st.session_state.get("slide_contents", [])
        if not slides:
            st.error("No slide contents found. Please regenerate the presentation in Tab 1.")
            return

        with st.spinner("Calling AI to design your presentation..."):
            try:
                prompt = build_beautify_prompt(slides)
                model_info = {"url": effective_url, "model": effective_model, "type": effective_type}
                response_text = call_llm(model_info, api_key, prompt)
            except requests.exceptions.Timeout:
                st.error("API request timed out. Please try again or use a faster model.")
                return
            except requests.exceptions.HTTPError as e:
                st.error(f"API returned an error: {e.response.status_code} – {e.response.text[:300]}")
                return
            except Exception as e:
                st.error(f"Error calling AI: {e}")
                return

        st.text_area("AI Response (raw)", response_text, height=150, key="raw_ai_resp")

        with st.spinner("Applying beautification..."):
            design = _extract_json(response_text)
            if not design:
                st.error("Could not parse AI response as JSON. The raw response is shown above. You can try again — sometimes the model returns extra text.")
                return

            try:
                beautified = apply_beautification(st.session_state.base_pptx, design, slides)
                st.session_state.beautified_pptx = beautified
                fn = st.session_state.get("base_filename", "presentation.pptx").replace(".pptx", "_beautified.pptx")
                st.session_state.beautified_filename = fn
                st.success("Presentation beautified successfully! 🎉")
            except Exception as e:
                st.error(f"Error applying beautification: {e}")

    if st.session_state.get("beautified_pptx"):
        st.download_button(
            label="⬇️ Download Beautified Presentation",
            data=st.session_state.beautified_pptx,
            file_name=st.session_state.get("beautified_filename", "presentation_beautified.pptx"),
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )

# ═══════════════════════════════════════════════
# UI: TAB 3 – SUPERUSER PANEL
# ═══════════════════════════════════════════════
def show_tab3():
    st.header("🛡️ Superuser Panel – User Management")

    users = get_all_users()
    if not users:
        st.info("No users registered yet.")
        return

    search = st.text_input("Search by access code", placeholder="KW-...", key="su_search")
    filtered = [u for u in users if search.upper() in u["code"].upper()] if search else users

    st.markdown(f"**{len(filtered)} user(s)** found")

    for u in filtered:
        code = u["code"]
        usage = u["usage_count"]
        blocked = bool(u["is_blocked"])
        expiry = u["access_expiry"]
        created = u["created_at"]
        last_login = u["last_login"]

        now = datetime.now()
        if expiry:
            exp_dt = datetime.fromisoformat(expiry)
            if exp_dt > now and not blocked:
                status = "🟢 Active Subscription"
                status_color = "green"
            elif exp_dt <= now:
                status = "🔴 Expired"
                status_color = "red"
            else:
                status = "🟡 Blocked (has subscription)"
                status_color = "orange"
        elif blocked:
            status = "🔴 Blocked"
            status_color = "red"
        elif usage >= FREE_USAGE_LIMIT:
            status = "🟡 Free uses exhausted"
            status_color = "orange"
        else:
            remaining = FREE_USAGE_LIMIT - usage
            status = f"🟢 Free user ({remaining} use{'s' if remaining != 1 else ''} left)"
            status_color = "green"

        with st.container():
            col1, col2, col3 = st.columns([3, 2, 2])
            with col1:
                st.markdown(f"**`{code}`**")
                st.caption(f"Created: {created[:19] if created else 'N/A'} | "
                           f"Last login: {last_login[:19] if last_login else 'N/A'}")
            with col2:
                st.markdown(f"Uses: **{usage}**/{FREE_USAGE_LIMIT}")
                st.markdown(f"<span style='color:{status_color}'>{status}</span>",
                            unsafe_allow_html=True)
                if expiry:
                    st.caption(f"Expiry: {expiry[:19]}")
            with col3:
                dur = st.selectbox("Duration", ["3 months", "6 months", "1 year"],
                                    key=f"dur_{code}", index=1,
                                    label_visibility="collapsed")
                months = {"3 months": 3, "6 months": 6, "1 year": 12}[dur]
                c_btn1, c_btn2 = st.columns(2)
                with c_btn1:
                    if st.button("✅ Unblock", key=f"unblk_{code}", use_container_width=True):
                        unblock_user(code, months)
                        st.success(f"Unblocked {code} for {dur}")
                        st.rerun()
                with c_btn2:
                    if st.button("🚫 Block", key=f"blk_{code}", use_container_width=True):
                        block_user(code)
                        st.warning(f"Blocked {code}")
                        st.rerun()
            st.divider()

# ═══════════════════════════════════════════════
# MAIN APPLICATION
# ═══════════════════════════════════════════════
def main():
    st.set_page_config(page_title="Knowledge Warehouse – PPT Generator",
                       page_icon="📄", layout="wide", initial_sidebar_state="collapsed")
    init_db()

    _defaults = {
        "logged_in": None,
        "user_code": None,
        "is_superuser": None,
        "base_pptx": None,
        "beautified_pptx": None,
        "slide_contents": None,
        "base_filename": None,
        "beautified_filename": None,
        "reg_code": None,
        "reg_email_msg": None,
        "reg_email_ok": None,
    }
    
    for key, default in _defaults.items():
        if key not in st.session_state:
            st.session_state[key] = default

    if not st.session_state.logged_in:
        show_landing()
        return

    st.markdown(CUSTOM_CSS, unsafe_allow_html=True)
    st.markdown("""<div class="kw-header" style="padding:1rem 2rem;">
        <h1 style="font-size:1.5rem;margin:0;">📄 AI Presentation Generator</h1>
        <p class="sub" style="margin:0;font-size:0.9rem;">by Knowledge Warehouse</p>
    </div>""", unsafe_allow_html=True)

    with st.sidebar:
        if st.session_state.is_superuser:
            st.success("🛡️ Superuser")
        else:
            st.info(f"👤 Code: `{st.session_state.user_code}`")
            user = get_user(st.session_state.user_code)
            if user:
                if user["access_expiry"]:
                    exp = datetime.fromisoformat(user["access_expiry"])
                    if exp > datetime.now():
                        st.success(f"Subscription active until {exp.strftime('%Y-%m-%d')}")
                    else:
                        st.warning("Subscription expired")
                else:
                    remaining = max(0, FREE_USAGE_LIMIT - user["usage_count"])
                    st.caption(f"Free uses remaining: {remaining}")

        if st.button("🚪 Logout", use_container_width=True):
            for key in list(st.session_state.keys()):
                del st.session_state[key]
            st.rerun()

    if st.session_state.is_superuser:
        tab1, tab2, tab3 = st.tabs(["📄 Generate Presentation",
                                     "🎨 Beautify with AI",
                                     "🛡️ Superuser Panel"])
        with tab1:
            show_tab1()
        with tab2:
            show_tab2()
        with tab3:
            show_tab3()
    else:
        tab1, tab2 = st.tabs(["📄 Generate Presentation",
                               "🎨 Beautify with AI"])
        with tab1:
            show_tab1()
        with tab2:
            show_tab2()

if __name__ == "__main__":
    main()
